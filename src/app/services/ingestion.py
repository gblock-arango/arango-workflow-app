"""Document parsing, semantic chunking, and duplicate detection.

Supports:

* **PDF**  via ``pymupdf`` (``fitz``)
* **DOCX** via ``python-docx``
* **PPTX** via ``python-pptx``
* **DOC**  (legacy Word binary, pre-2007) via a LibreOffice subprocess
  that converts ``.doc`` -> ``.docx`` and then reuses :func:`parse_docx`.
  Requires LibreOffice/soffice on PATH (``brew install --cask libreoffice``
  on macOS, ``apt install libreoffice-core`` on Debian). Fails loudly
  with a clear error if not present.
* **Markdown** plain-text parser
"""

from __future__ import annotations

import hashlib
import io
import logging
import os
import re
import shutil
import subprocess
import tempfile
from collections.abc import Iterable
from dataclasses import dataclass, field
from typing import Any

import fitz
import tiktoken
from docx import Document
from pptx import Presentation

log = logging.getLogger(__name__)

# Subprocess timeout for the LibreOffice .doc -> .docx conversion.
# 60s handles large legacy decks; the convert pass itself rarely exceeds 5s.
_LIBREOFFICE_TIMEOUT_SECONDS = 60

_DEFAULT_MAX_TOKENS = 512
_TIKTOKEN_MODEL = "cl100k_base"


@dataclass
class Section:
    heading: str
    text: str
    page_number: int | None = None


@dataclass
class ParsedDocument:
    sections: list[Section] = field(default_factory=list)
    title: str = ""
    author: str = ""
    page_count: int = 0


@dataclass
class Chunk:
    text: str
    chunk_index: int
    source_page: int | None
    section_heading: str
    token_count: int


def compute_file_hash(content: bytes) -> str:
    """SHA-256 hex digest of raw file bytes."""
    return hashlib.sha256(content).hexdigest()


def _token_count(text: str) -> int:
    enc = tiktoken.get_encoding(_TIKTOKEN_MODEL)
    return len(enc.encode(text))


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------


def parse_pdf(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a PDF using pymupdf (fitz)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    parsed = ParsedDocument(page_count=len(doc))

    metadata = doc.metadata or {}
    parsed.title = metadata.get("title", "") or ""
    parsed.author = metadata.get("author", "") or ""

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        current_heading = ""
        current_text_parts: list[str] = []

        for block in blocks:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                line_text = "".join(span["text"] for span in line.get("spans", []))
                if not line_text.strip():
                    continue

                max_size = max(
                    (span.get("size", 0) for span in line.get("spans", [])),
                    default=0,
                )
                is_bold = any(
                    "bold" in (span.get("font", "").lower()) for span in line.get("spans", [])
                )

                if max_size >= 14 or (is_bold and max_size >= 12):
                    if current_text_parts:
                        parsed.sections.append(
                            Section(
                                heading=current_heading,
                                text="\n".join(current_text_parts).strip(),
                                page_number=page_num,
                            )
                        )
                        current_text_parts = []
                    current_heading = line_text.strip()
                else:
                    current_text_parts.append(line_text)

        if current_text_parts:
            parsed.sections.append(
                Section(
                    heading=current_heading,
                    text="\n".join(current_text_parts).strip(),
                    page_number=page_num,
                )
            )

    doc.close()

    if not parsed.sections:
        full_text = ""
        reopened = fitz.open(stream=file_bytes, filetype="pdf")
        for page in reopened:
            full_text += page.get_text() + "\n"
        reopened.close()
        if full_text.strip():
            parsed.sections.append(Section(heading="", text=full_text.strip(), page_number=1))

    return parsed


def parse_docx(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a DOCX using python-docx."""
    doc = Document(io.BytesIO(file_bytes))

    parsed = ParsedDocument()
    core = doc.core_properties
    parsed.title = core.title or ""
    parsed.author = core.author or ""

    current_heading = ""
    current_text_parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style and para.style.name and para.style.name.startswith("Heading"):
            if current_text_parts:
                parsed.sections.append(
                    Section(heading=current_heading, text="\n".join(current_text_parts))
                )
                current_text_parts = []
            current_heading = text
        else:
            current_text_parts.append(text)

    if current_text_parts:
        parsed.sections.append(Section(heading=current_heading, text="\n".join(current_text_parts)))

    return parsed


def parse_pptx(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a PPTX (PowerPoint) using python-pptx.

    Each slide becomes one :class:`Section` so downstream chunking and
    provenance line up with what users intuitively expect ("page 7 of
    that deck"):

    * ``page_number`` = 1-based slide index (matches PDF semantics).
    * ``heading`` = slide title placeholder if present, else "".
    * ``text`` = concatenation of every other text-bearing shape on the
      slide (body placeholders, text boxes, table cells, grouped shapes
      flattened recursively). Notes pages are appended below the slide
      body, prefixed with ``[Notes]`` so the LLM can distinguish them.

    Empty slides (no extractable text and no notes) are dropped to keep
    chunk counts honest.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    parsed = ParsedDocument(page_count=len(prs.slides))

    core = prs.core_properties
    parsed.title = (core.title or "") if core else ""
    parsed.author = (core.author or "") if core else ""

    for slide_index, slide in enumerate(prs.slides, start=1):
        heading = _pptx_slide_title(slide)
        body_parts = _pptx_collect_text(slide.shapes, exclude_title=True)

        # Speaker notes: powerful provenance signal, often where the
        # actual narrative lives in survey decks.
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
                notes_text = notes.strip()
        except Exception:
            # Don't let a malformed notes pane kill the whole parse.
            log.warning("pptx parse: failed to read notes for slide %d", slide_index, exc_info=True)

        if notes_text:
            body_parts.append(f"[Notes] {notes_text}")

        body = "\n".join(p for p in body_parts if p).strip()
        if not body and not heading:
            continue

        parsed.sections.append(Section(heading=heading, text=body, page_number=slide_index))

    return parsed


def _pptx_slide_title(slide: object) -> str:
    """Return the slide's title-placeholder text if any, else ""."""
    try:
        title_shape = slide.shapes.title  # type: ignore[attr-defined]
    except Exception:
        return ""
    if title_shape is None:
        return ""
    text = getattr(title_shape, "text", "") or ""
    return text.strip()


def _pptx_collect_text(
    shapes: Iterable[Any],
    *,
    exclude_title: bool,
) -> list[str]:
    """Walk a python-pptx shape collection (incl. groups + tables) for text.

    ``exclude_title`` skips the slide title shape so we don't double-count
    it (the caller already extracted it as the section heading).
    """
    out: list[str] = []
    for shape in shapes:
        if exclude_title and getattr(shape, "is_placeholder", False):
            ph = getattr(shape, "placeholder_format", None)
            # Title placeholder idx is 0 in OpenXML.
            if ph is not None and getattr(ph, "idx", None) == 0:
                continue

        # Grouped shapes -> recurse.
        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            out.extend(_pptx_collect_text(shape.shapes, exclude_title=False))
            continue

        # Tables: row by row, cell by cell.
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join((cell.text or "").strip() for cell in row.cells).strip(" |")
                if row_text:
                    out.append(row_text)
            continue

        # Plain text frames.
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()
            if text:
                out.append(text)
    return out


def parse_doc(file_bytes: bytes) -> ParsedDocument:
    """Extract sections from a legacy ``.doc`` (Word 97-2003 binary) file.

    Strategy: shell out to LibreOffice in headless mode to convert the
    binary ``.doc`` to ``.docx``, then reuse :func:`parse_docx` so we
    keep one code path for Word styles + heading detection.

    LibreOffice / ``soffice`` must be on PATH. We probe at call-time
    rather than import-time so the rest of the module still imports
    on a host that has not yet installed it.

    Raises
    ------
    RuntimeError
        When LibreOffice is not installed or the conversion fails. The
        error message tells the operator exactly what to install.
    """
    soffice = _find_libreoffice()
    if soffice is None:
        raise RuntimeError(
            "Cannot parse legacy .doc files: LibreOffice (soffice) is not "
            "installed. Install it with `brew install --cask libreoffice` "
            "(macOS) or `apt install libreoffice-core` (Debian/Ubuntu) and "
            "retry. Alternatively, convert the file to .docx and re-upload."
        )

    with tempfile.TemporaryDirectory(prefix="aoe_doc_") as tmpdir:
        in_path = os.path.join(tmpdir, "in.doc")
        out_dir = os.path.join(tmpdir, "out")
        os.makedirs(out_dir, exist_ok=True)
        with open(in_path, "wb") as fh:
            fh.write(file_bytes)

        try:
            result = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "docx",
                    "--outdir",
                    out_dir,
                    in_path,
                ],
                capture_output=True,
                timeout=_LIBREOFFICE_TIMEOUT_SECONDS,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError(
                f"LibreOffice .doc -> .docx conversion timed out after "
                f"{_LIBREOFFICE_TIMEOUT_SECONDS}s. The file may be malformed "
                f"or unusually large."
            ) from exc

        if result.returncode != 0:
            stderr = (result.stderr or b"").decode("utf-8", errors="replace").strip()
            raise RuntimeError(
                f"LibreOffice conversion failed (exit={result.returncode}): "
                f"{stderr or '<no stderr>'}"
            )

        # The converted file lands in out_dir as <basename>.docx.
        candidates = [f for f in os.listdir(out_dir) if f.endswith(".docx")]
        if not candidates:
            raise RuntimeError(
                "LibreOffice conversion produced no .docx output. Stderr: "
                + (result.stderr or b"").decode("utf-8", errors="replace").strip()
            )

        with open(os.path.join(out_dir, candidates[0]), "rb") as fh:
            converted_bytes = fh.read()

    return parse_docx(converted_bytes)


def _find_libreoffice() -> str | None:
    """Locate the LibreOffice headless binary, or return None if absent.

    Checks PATH for ``soffice`` and ``libreoffice``, then the standard
    macOS install location. Order matters: PATH first so an admin can
    pin a specific version.
    """
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    mac_default = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_default) and os.access(mac_default, os.X_OK):
        return mac_default
    return None


_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)", re.MULTILINE)


def parse_markdown(text: str) -> ParsedDocument:
    """Extract sections from Markdown text using heading boundaries."""
    parsed = ParsedDocument()

    lines = text.split("\n")
    first_heading = ""
    for line in lines:
        m = _MD_HEADING_RE.match(line)
        if m:
            first_heading = m.group(2).strip()
            break
    parsed.title = first_heading

    current_heading = ""
    current_text_parts: list[str] = []

    for line in lines:
        m = _MD_HEADING_RE.match(line)
        if m:
            if current_text_parts:
                parsed.sections.append(
                    Section(heading=current_heading, text="\n".join(current_text_parts).strip())
                )
                current_text_parts = []
            current_heading = m.group(2).strip()
        else:
            current_text_parts.append(line)

    if current_text_parts:
        body = "\n".join(current_text_parts).strip()
        if body:
            parsed.sections.append(Section(heading=current_heading, text=body))

    return parsed


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------


def _split_into_paragraphs(text: str) -> list[str]:
    """Split text at blank-line boundaries, keeping non-empty paragraphs."""
    paragraphs = re.split(r"\n\s*\n", text)
    return [p.strip() for p in paragraphs if p.strip()]


def chunk_document(
    parsed: ParsedDocument,
    max_tokens: int = _DEFAULT_MAX_TOKENS,
) -> list[Chunk]:
    """Chunk a parsed document at section / paragraph boundaries.

    Each chunk respects ``max_tokens`` (counted via tiktoken ``cl100k_base``).
    Chunks preserve source page and section heading metadata.
    """
    chunks: list[Chunk] = []
    idx = 0

    for section in parsed.sections:
        paragraphs = _split_into_paragraphs(section.text)
        if not paragraphs:
            continue

        current_parts: list[str] = []
        current_tokens = 0

        for para in paragraphs:
            para_tokens = _token_count(para)

            if para_tokens > max_tokens:
                if current_parts:
                    merged = "\n\n".join(current_parts)
                    chunks.append(
                        Chunk(
                            text=merged,
                            chunk_index=idx,
                            source_page=section.page_number,
                            section_heading=section.heading,
                            token_count=_token_count(merged),
                        )
                    )
                    idx += 1
                    current_parts = []
                    current_tokens = 0

                words = para.split()
                sub_parts: list[str] = []
                sub_tokens = 0
                for word in words:
                    word_tokens = _token_count(word + " ")
                    if sub_tokens + word_tokens > max_tokens and sub_parts:
                        sub_text = " ".join(sub_parts)
                        chunks.append(
                            Chunk(
                                text=sub_text,
                                chunk_index=idx,
                                source_page=section.page_number,
                                section_heading=section.heading,
                                token_count=_token_count(sub_text),
                            )
                        )
                        idx += 1
                        sub_parts = []
                        sub_tokens = 0
                    sub_parts.append(word)
                    sub_tokens += word_tokens

                if sub_parts:
                    sub_text = " ".join(sub_parts)
                    chunks.append(
                        Chunk(
                            text=sub_text,
                            chunk_index=idx,
                            source_page=section.page_number,
                            section_heading=section.heading,
                            token_count=_token_count(sub_text),
                        )
                    )
                    idx += 1
                continue

            if current_tokens + para_tokens > max_tokens and current_parts:
                merged = "\n\n".join(current_parts)
                chunks.append(
                    Chunk(
                        text=merged,
                        chunk_index=idx,
                        source_page=section.page_number,
                        section_heading=section.heading,
                        token_count=_token_count(merged),
                    )
                )
                idx += 1
                current_parts = []
                current_tokens = 0

            current_parts.append(para)
            current_tokens += para_tokens

        if current_parts:
            merged = "\n\n".join(current_parts)
            chunks.append(
                Chunk(
                    text=merged,
                    chunk_index=idx,
                    source_page=section.page_number,
                    section_heading=section.heading,
                    token_count=_token_count(merged),
                )
            )
            idx += 1

    return chunks
