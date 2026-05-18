"""Unit tests for app.services.ingestion — parsing, chunking, hashing."""

from __future__ import annotations

from unittest.mock import MagicMock, patch

from app.services.ingestion import (
    ParsedDocument,
    Section,
    chunk_document,
    compute_file_hash,
    parse_markdown,
)

# ---------------------------------------------------------------------------
# compute_file_hash
# ---------------------------------------------------------------------------


class TestComputeFileHash:
    def test_deterministic(self):
        content = b"hello world"
        h1 = compute_file_hash(content)
        h2 = compute_file_hash(content)
        assert h1 == h2

    def test_different_content_different_hash(self):
        assert compute_file_hash(b"a") != compute_file_hash(b"b")

    def test_empty_bytes(self):
        h = compute_file_hash(b"")
        assert isinstance(h, str)
        assert len(h) == 64

    def test_sha256_known_value(self):
        expected = "9f86d081884c7d659a2feaa0c55ad015a3bf4f1b2b0b822cd15d6c15b0f00a08"
        assert compute_file_hash(b"test") == expected


# ---------------------------------------------------------------------------
# parse_markdown
# ---------------------------------------------------------------------------


class TestParseMarkdown:
    def test_single_section(self):
        md = "# Title\n\nSome body text here."
        parsed = parse_markdown(md)
        assert parsed.title == "Title"
        assert len(parsed.sections) == 1
        assert parsed.sections[0].heading == "Title"
        assert "body text" in parsed.sections[0].text

    def test_multiple_sections(self):
        md = "# Intro\n\nIntro text.\n\n## Details\n\nDetail text."
        parsed = parse_markdown(md)
        assert len(parsed.sections) == 2
        assert parsed.sections[0].heading == "Intro"
        assert parsed.sections[1].heading == "Details"

    def test_empty_markdown(self):
        parsed = parse_markdown("")
        assert parsed.sections == []
        assert parsed.title == ""

    def test_no_headings(self):
        md = "Just a paragraph of text without any headings."
        parsed = parse_markdown(md)
        assert len(parsed.sections) == 1
        assert parsed.sections[0].heading == ""

    def test_heading_levels(self):
        md = "# H1\n\nh1 text\n\n### H3\n\nh3 text"
        parsed = parse_markdown(md)
        assert parsed.sections[0].heading == "H1"
        assert parsed.sections[1].heading == "H3"


# ---------------------------------------------------------------------------
# parse_pdf (mocked)
# ---------------------------------------------------------------------------


class TestParsePdf:
    @patch("app.services.ingestion.fitz")
    def test_basic_extraction(self, mock_fitz: MagicMock):
        mock_page = MagicMock()
        mock_page.get_text.return_value = {
            "blocks": [
                {
                    "type": 0,
                    "lines": [
                        {"spans": [{"text": "A normal paragraph.", "size": 11, "font": "Regular"}]}
                    ],
                }
            ]
        }

        mock_doc = MagicMock()
        mock_doc.__iter__ = MagicMock(return_value=iter([mock_page]))
        mock_doc.__len__ = MagicMock(return_value=1)
        mock_doc.metadata = {"title": "Test", "author": "Auth"}

        mock_fitz.open.return_value = mock_doc
        mock_fitz.TEXT_PRESERVE_WHITESPACE = 1

        from app.services.ingestion import parse_pdf

        parsed = parse_pdf(b"fake-pdf-bytes")

        assert parsed.title == "Test"
        assert parsed.author == "Auth"
        assert len(parsed.sections) >= 1

    @patch("app.services.ingestion.fitz")
    def test_empty_pdf(self, mock_fitz: MagicMock):
        mock_doc = MagicMock()
        mock_doc.__iter__ = MagicMock(return_value=iter([]))
        mock_doc.__len__ = MagicMock(return_value=0)
        mock_doc.metadata = {}

        reopened_doc = MagicMock()
        reopened_page = MagicMock()
        reopened_page.get_text.return_value = ""
        reopened_doc.__iter__ = MagicMock(return_value=iter([reopened_page]))

        mock_fitz.open.side_effect = [mock_doc, reopened_doc]
        mock_fitz.TEXT_PRESERVE_WHITESPACE = 1

        from app.services.ingestion import parse_pdf

        parsed = parse_pdf(b"fake-empty-pdf")
        assert parsed.page_count == 0


# ---------------------------------------------------------------------------
# parse_docx (mocked)
# ---------------------------------------------------------------------------


class TestParseDocx:
    @patch("app.services.ingestion.Document")
    def test_basic_extraction(self, mock_document_cls: MagicMock):
        mock_heading_style = MagicMock()
        mock_heading_style.name = "Heading 1"

        mock_normal_style = MagicMock()
        mock_normal_style.name = "Normal"

        heading_para = MagicMock()
        heading_para.text = "My Heading"
        heading_para.style = mock_heading_style

        body_para = MagicMock()
        body_para.text = "Some body content here."
        body_para.style = mock_normal_style

        mock_core = MagicMock()
        mock_core.title = "Doc Title"
        mock_core.author = "Author Name"

        mock_doc = MagicMock()
        mock_doc.paragraphs = [heading_para, body_para]
        mock_doc.core_properties = mock_core
        mock_document_cls.return_value = mock_doc

        from app.services.ingestion import parse_docx

        parsed = parse_docx(b"fake-docx")

        assert parsed.title == "Doc Title"
        assert parsed.author == "Author Name"
        assert len(parsed.sections) == 1
        assert parsed.sections[0].heading == "My Heading"
        assert "body content" in parsed.sections[0].text


# ---------------------------------------------------------------------------
# chunk_document (tiktoken mocked)
# ---------------------------------------------------------------------------


def _fake_token_count(text: str) -> int:
    """Approximate token count for tests: ~1 token per 4 chars."""
    return max(1, len(text) // 4)


@patch("app.services.ingestion._token_count", side_effect=_fake_token_count)
class TestChunkDocument:
    def test_single_short_section(self, _mock_tc: MagicMock):
        parsed = ParsedDocument(
            sections=[Section(heading="Intro", text="Short text.", page_number=1)]
        )
        chunks = chunk_document(parsed, max_tokens=512)
        assert len(chunks) == 1
        assert chunks[0].chunk_index == 0
        assert chunks[0].section_heading == "Intro"
        assert chunks[0].source_page == 1

    def test_respects_max_tokens(self, _mock_tc: MagicMock):
        long_text = " ".join(["word"] * 2000)
        parsed = ParsedDocument(sections=[Section(heading="Long", text=long_text, page_number=1)])
        chunks = chunk_document(parsed, max_tokens=50)
        assert len(chunks) > 1

    def test_empty_document(self, _mock_tc: MagicMock):
        parsed = ParsedDocument(sections=[])
        chunks = chunk_document(parsed)
        assert chunks == []

    def test_multiple_sections(self, _mock_tc: MagicMock):
        parsed = ParsedDocument(
            sections=[
                Section(heading="A", text="Text A", page_number=1),
                Section(heading="B", text="Text B", page_number=2),
            ]
        )
        chunks = chunk_document(parsed)
        assert len(chunks) == 2
        assert chunks[0].section_heading == "A"
        assert chunks[1].section_heading == "B"

    def test_chunk_indexes_sequential(self, _mock_tc: MagicMock):
        parsed = ParsedDocument(
            sections=[
                Section(heading="A", text="Text A", page_number=1),
                Section(heading="B", text="Text B", page_number=2),
            ]
        )
        chunks = chunk_document(parsed)
        for i, chunk in enumerate(chunks):
            assert chunk.chunk_index == i

    def test_paragraph_boundary_splitting(self, _mock_tc: MagicMock):
        text = (
            "First paragraph with enough words to matter.\n\n"
            "Second paragraph also with content.\n\n"
            "Third paragraph here too."
        )
        parsed = ParsedDocument(sections=[Section(heading="S", text=text, page_number=1)])
        chunks = chunk_document(parsed, max_tokens=15)
        assert len(chunks) >= 2


# ---------------------------------------------------------------------------
# parse_pptx (real round-trip with python-pptx)
# ---------------------------------------------------------------------------


def _build_synthetic_pptx() -> bytes:
    """Create an in-memory .pptx with three slides for round-trip tests."""
    import io as _io

    from pptx import Presentation as _Pres
    from pptx.util import Inches as _In

    prs = _Pres()

    # Slide 1: title + bullets layout, with notes.
    layout = prs.slide_layouts[1]  # "Title and Content"
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Healthcare Survey 2024"
    body = s1.placeholders[1]
    body.text = "Key finding: claim costs up 8%"
    p2 = body.text_frame.add_paragraph()
    p2.text = "Mental health utilisation up 15%"
    s1.notes_slide.notes_text_frame.text = "Speaker note: emphasise mental-health spike"

    # Slide 2: blank layout + a freeform text box + a table.
    blank = prs.slide_layouts[5]  # "Title Only"
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "Benefit Preferences"
    tx = s2.shapes.add_textbox(_In(1), _In(2), _In(5), _In(2))
    tx.text_frame.text = "Most-valued benefit: 401(k) match"
    rows, cols = 2, 2
    table = s2.shapes.add_table(rows, cols, _In(1), _In(4), _In(5), _In(1.5)).table
    table.cell(0, 0).text = "Benefit"
    table.cell(0, 1).text = "Rank"
    table.cell(1, 0).text = "Health insurance"
    table.cell(1, 1).text = "1"

    # Slide 3: deliberately empty (Blank layout has no title placeholder
    # and we add no shapes) -> should be skipped by parse_pptx.
    prs.slides.add_slide(prs.slide_layouts[6])  # "Blank"

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestParsePptx:
    """Round-trip tests using python-pptx to build deterministic fixtures."""

    def test_parses_real_pptx_into_sections_per_slide(self):
        from app.services.ingestion import parse_pptx

        pptx_bytes = _build_synthetic_pptx()
        parsed = parse_pptx(pptx_bytes)

        # Slide 3 is dropped (empty), so only 2 sections.
        assert len(parsed.sections) == 2
        # page_count reflects raw slide count, not filtered sections.
        assert parsed.page_count == 3

    def test_first_slide_title_and_body_extracted(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(_build_synthetic_pptx())
        s1 = parsed.sections[0]
        assert s1.heading == "Healthcare Survey 2024"
        assert s1.page_number == 1
        assert "claim costs up 8%" in s1.text
        assert "Mental health utilisation up 15%" in s1.text

    def test_speaker_notes_appended_with_marker(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(_build_synthetic_pptx())
        assert "[Notes]" in parsed.sections[0].text
        assert "mental-health spike" in parsed.sections[0].text

    def test_title_not_double_counted_in_body(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(_build_synthetic_pptx())
        # The slide title appears exactly once -- as the heading, not also in body.
        assert parsed.sections[0].heading.count("Healthcare Survey 2024") == 1
        assert "Healthcare Survey 2024" not in parsed.sections[0].text

    def test_text_boxes_and_tables_extracted(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(_build_synthetic_pptx())
        s2 = parsed.sections[1]
        assert s2.heading == "Benefit Preferences"
        assert "Most-valued benefit: 401(k) match" in s2.text
        # Table rows joined as `cell | cell`.
        assert "Benefit | Rank" in s2.text
        assert "Health insurance | 1" in s2.text

    def test_empty_slide_is_dropped(self):
        from app.services.ingestion import parse_pptx

        parsed = parse_pptx(_build_synthetic_pptx())
        # Slide 3 produced no heading and no body -> not in sections.
        page_numbers = {s.page_number for s in parsed.sections}
        assert 3 not in page_numbers

    def test_invalid_bytes_raises(self):
        """Garbage bytes -> the underlying zipfile layer rejects it.

        python-pptx wraps a ZIP container, so non-ZIP input surfaces as
        ``zipfile.BadZipFile``. We accept any subclass of OSError /
        zipfile.BadZipFile to stay tolerant of future python-pptx
        version changes that may wrap this in their own exception.
        """
        import zipfile

        import pytest as _pytest

        from app.services.ingestion import parse_pptx

        with _pytest.raises((zipfile.BadZipFile, OSError, ValueError)):
            parse_pptx(b"definitely not a real pptx")


# ---------------------------------------------------------------------------
# parse_doc (mocked LibreOffice subprocess + real parse_docx on output)
# ---------------------------------------------------------------------------


class TestParseDoc:
    def test_raises_when_libreoffice_missing(self):
        import pytest as _pytest

        from app.services import ingestion as _ing

        with (
            patch.object(_ing, "_find_libreoffice", return_value=None),
            _pytest.raises(RuntimeError, match="LibreOffice"),
        ):
            _ing.parse_doc(b"fake-doc-bytes")

    def test_calls_soffice_with_expected_args(self, tmp_path):
        # Build a real .docx (one paragraph) that we'll pretend soffice produced.
        from docx import Document as _Document

        from app.services import ingestion as _ing

        d = _Document()
        d.add_heading("Converted Heading", level=1)
        d.add_paragraph("Converted body text.")
        import io as _io

        buf = _io.BytesIO()
        d.save(buf)
        docx_bytes = buf.getvalue()

        captured: dict[str, list[str]] = {}

        def _fake_run(cmd, **kwargs):
            captured["cmd"] = list(cmd)
            # Mimic LibreOffice: drop a .docx into the --outdir.
            outdir_idx = cmd.index("--outdir") + 1
            outdir = cmd[outdir_idx]
            with open(f"{outdir}/in.docx", "wb") as fh:
                fh.write(docx_bytes)
            result = MagicMock()
            result.returncode = 0
            result.stderr = b""
            return result

        with (
            patch.object(_ing, "_find_libreoffice", return_value="/fake/soffice"),
            patch("app.services.ingestion.subprocess.run", side_effect=_fake_run),
        ):
            parsed = _ing.parse_doc(b"fake-doc-bytes")

        assert captured["cmd"][0] == "/fake/soffice"
        assert "--headless" in captured["cmd"]
        assert "--convert-to" in captured["cmd"]
        assert captured["cmd"][captured["cmd"].index("--convert-to") + 1] == "docx"
        # parse_docx round-tripped the converted output.
        assert parsed.sections[0].heading == "Converted Heading"
        assert "Converted body text" in parsed.sections[0].text

    def test_subprocess_failure_raises_with_stderr(self):
        import pytest as _pytest

        from app.services import ingestion as _ing

        result = MagicMock()
        result.returncode = 2
        result.stderr = b"soffice: source file not found"

        with (
            patch.object(_ing, "_find_libreoffice", return_value="/fake/soffice"),
            patch("app.services.ingestion.subprocess.run", return_value=result),
            _pytest.raises(RuntimeError, match="exit=2"),
        ):
            _ing.parse_doc(b"fake")

    def test_timeout_raises_with_clear_message(self):
        import subprocess as _sp

        import pytest as _pytest

        from app.services import ingestion as _ing

        with (
            patch.object(_ing, "_find_libreoffice", return_value="/fake/soffice"),
            patch(
                "app.services.ingestion.subprocess.run",
                side_effect=_sp.TimeoutExpired(cmd="soffice", timeout=60),
            ),
            _pytest.raises(RuntimeError, match="timed out"),
        ):
            _ing.parse_doc(b"fake")

    def test_no_docx_output_raises(self):
        import pytest as _pytest

        from app.services import ingestion as _ing

        result = MagicMock()
        result.returncode = 0
        result.stderr = b""

        # subprocess returns success but produces no .docx (LO gone weird)
        with (
            patch.object(_ing, "_find_libreoffice", return_value="/fake/soffice"),
            patch("app.services.ingestion.subprocess.run", return_value=result),
            _pytest.raises(RuntimeError, match=r"no \.docx output"),
        ):
            _ing.parse_doc(b"fake")


# ---------------------------------------------------------------------------
# _find_libreoffice
# ---------------------------------------------------------------------------


class TestFindLibreoffice:
    def test_prefers_path_soffice(self):
        from app.services import ingestion as _ing

        with patch("app.services.ingestion.shutil.which") as which:
            which.side_effect = lambda name: "/usr/local/bin/soffice" if name == "soffice" else None
            assert _ing._find_libreoffice() == "/usr/local/bin/soffice"

    def test_falls_back_to_libreoffice_name(self):
        from app.services import ingestion as _ing

        with patch("app.services.ingestion.shutil.which") as which:
            which.side_effect = lambda name: (
                "/usr/bin/libreoffice" if name == "libreoffice" else None
            )
            assert _ing._find_libreoffice() == "/usr/bin/libreoffice"

    def test_falls_back_to_mac_default_when_path_empty(self):
        from app.services import ingestion as _ing

        mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        with (
            patch("app.services.ingestion.shutil.which", return_value=None),
            patch("app.services.ingestion.os.path.isfile", return_value=True),
            patch("app.services.ingestion.os.access", return_value=True),
        ):
            assert _ing._find_libreoffice() == mac_path

    def test_returns_none_when_truly_absent(self):
        from app.services import ingestion as _ing

        with (
            patch("app.services.ingestion.shutil.which", return_value=None),
            patch("app.services.ingestion.os.path.isfile", return_value=False),
        ):
            assert _ing._find_libreoffice() is None
